"""
engines/document.py — FileZaar Document Engine v3
- Preserves images in all conversions (PDF/DOCX/PPTX/HTML)
- Unicode / multi-language support (Hindi, Arabic, CJK, etc.)
- Original filename preserved via output_stem
- Added: EPUB, XML, PPTX creation conversions
"""
import asyncio, shutil, html as _html, csv as _csv, io, re, base64
from pathlib import Path
from typing import Callable, Awaitable, Optional
from core.logger import get_logger
logger = get_logger(__name__)

_PANDOC = None
def _find_pandoc():
    global _PANDOC
    if _PANDOC is None:
        _PANDOC = shutil.which("pandoc") or ""
    return _PANDOC or None

def _sanitize(text: str) -> str:
    """Remove NULL/control chars but keep all Unicode (Hindi, CJK, Arabic, etc.)"""
    if not text: return ""
    cleaned = []
    for ch in text:
        cp = ord(ch)
        if (cp in (0x09, 0x0A, 0x0D) or
            (0x20 <= cp <= 0xD7FF) or
            (0xE000 <= cp <= 0xFFFD) or
            (0x10000 <= cp <= 0x10FFFF)):
            cleaned.append(ch)
        else:
            cleaned.append(' ')
    return ''.join(cleaned)

def _img_to_b64(img_bytes: bytes, mime: str = "image/png") -> str:
    return f"data:{mime};base64,{base64.b64encode(img_bytes).decode()}"


class DocumentEngine:
    async def convert(self, input_path: Path, target_format: str,
                      output_dir: Path, progress_cb: Callable[[int,str],Awaitable[None]],
                      output_stem: Optional[str] = None) -> dict:
        src  = input_path.suffix.lstrip(".").lower()
        tgt  = target_format.lower().lstrip(".")
        inp  = input_path.resolve()
        odir = output_dir.resolve()
        odir.mkdir(parents=True, exist_ok=True)
        stem = output_stem or inp.stem
        await progress_cb(10, f"Converting {src.upper()} → {tgt.upper()}…")
        h = self._route(src, tgt)
        return await h(inp, tgt, odir, progress_cb, stem)

    # ── Route table ──────────────────────────────────────────────────────────
    def _route(self, src, tgt):
        table = {
            # MD
            ("md","html"):  self._md_html,   ("md","pdf"):  self._md_pdf,
            ("md","docx"):  self._pan,        ("md","txt"):  self._md_txt,
            ("md","epub"):  self._pan,        ("md","rtf"):  self._pan,
            ("md","csv"):   self._txt_csv,
            # TXT
            ("txt","html"): self._txt_html,   ("txt","pdf"): self._txt_pdf,
            ("txt","docx"): self._pan,        ("txt","epub"):self._pan,
            ("txt","xml"):  self._txt_xml,    ("txt","pptx"):self._txt_pptx,
            ("txt","csv"):  self._txt_csv,
            # HTML
            ("html","pdf"): self._html_pdf,   ("htm","pdf"): self._html_pdf,
            ("html","txt"): self._html_txt,   ("htm","txt"): self._html_txt,
            ("html","md"):  self._pan,        ("html","docx"):self._pan,
            ("html","epub"):self._pan,        ("html","rtf"):self._pan,
            ("html","xml"): self._html_xml,   ("html","csv"): self._html_csv,
            # DOCX
            ("docx","txt"): self._docx_txt,   ("docx","html"):self._docx_html,
            ("docx","pdf"): self._docx_pdf,   ("docx","md"):  self._pan,
            ("docx","epub"):self._pan,        ("docx","odt"): self._pan,
            ("docx","rtf"): self._pan,        ("docx","pptx"):self._docx_pptx,
            ("docx","csv"): self._docx_csv,   ("docx","xlsx"):self._docx_xlsx,
            # DOC
            ("doc","docx"): self._pan,        ("doc","pdf"):  self._pan,
            ("doc","txt"):  self._pan,        ("doc","html"): self._pan,
            # PDF
            ("pdf","txt"):  self._pdf_txt,    ("pdf","html"): self._pdf_html,
            ("pdf","docx"): self._pdf_docx,   ("pdf","md"):   self._pdf_md,
            ("pdf","xml"):  self._pdf_xml,
            ("pdf","epub"): self._pdf_via_txt, ("pdf","odt"):  self._pdf_via_txt,
            ("pdf","rtf"):  self._pdf_via_txt, ("pdf","csv"):  self._pdf_via_txt,
            ("pdf","json"): self._pdf_via_txt, ("pdf","xlsx"): self._pdf_via_txt,
            ("pdf","pptx"): self._pdf_via_txt,
            # PPTX
            # PPT (old binary format) - use LibreOffice
            ("ppt","pptx"): self._ppt_convert,  ("ppt","pdf"):  self._ppt_convert,
            ("ppt","txt"):  self._ppt_convert,   ("ppt","html"): self._ppt_convert,
            ("ppt","docx"): self._ppt_convert,
            ("pptx","txt"): self._pptx_txt,   ("pptx","html"):self._pptx_html,
            ("pptx","pdf"): self._pptx_pdf,   ("pptx","docx"): self._pptx_docx,
            ("pptx","md"):  self._pptx_via_txt,("pptx","odt"): self._pptx_via_txt,
            ("pptx","rtf"): self._pptx_via_txt,("pptx","epub"):self._pptx_via_txt,
            ("pptx","csv"): self._pptx_via_txt,("pptx","json"):self._pptx_via_txt,
            ("pptx","xlsx"):self._pptx_via_txt,("pptx","xml"): self._pptx_via_txt,
            # XLSX
            ("xlsx","csv"): self._xlsx_csv,   ("xlsx","txt"): self._xlsx_txt,
            ("xlsx","html"):self._xlsx_html,  ("xlsx","pdf"): self._xlsx_pdf,
            ("xlsx","json"):self._xlsx_json,  ("xlsx","xml"): self._xlsx_xml,
            ("xlsx","docx"):self._xlsx_via_txt,("xlsx","epub"):self._xlsx_via_txt,
            ("xlsx","md"):  self._xlsx_via_txt,("xlsx","odt"): self._xlsx_via_txt,
            ("xlsx","rtf"): self._xlsx_via_txt,("xlsx","pptx"):self._xlsx_via_txt,
            # CSV
            ("csv","xlsx"): self._csv_xlsx,   ("csv","txt"):  self._csv_txt,
            ("csv","html"): self._csv_html,   ("csv","pdf"):  self._csv_pdf,
            ("csv","json"): self._csv_json,   ("csv","xml"):  self._csv_xml,
            ("csv","docx"): self._csv_docx,   ("csv","md"):   self._csv_md,
            # JSON
            ("json","csv"): self._json_csv,   ("json","txt"): self._json_txt,
            ("json","html"):self._json_html,  ("json","xlsx"):self._json_xlsx,
            ("json","xml"): self._json_xml,   ("json","pdf"):  self._doc_via_html_pdf,
            # XML
            ("xml","json"): self._xml_json,   ("xml","txt"):  self._xml_txt,
            ("xml","html"): self._xml_html,   ("xml","pdf"):  self._xml_pdf,
            ("xml","csv"):  self._txt_csv,
            # Fix RC-3: explicit xlsx routes → openpyxl real XLSX
            ("md","xlsx"):   self._txt_xlsx,  ("txt","xlsx"):  self._txt_xlsx,
            ("html","xlsx"): self._txt_xlsx,  ("xml","xlsx"):  self._txt_xlsx,
            ("odt","xlsx"):  self._doc_via_txt_xlsx, ("rtf","xlsx"):  self._txt_xlsx,
            ("epub","xlsx"): self._doc_via_txt_xlsx,
            # CSS
            ("css","txt"):  self._copy_txt,   ("css","html"): self._css_html,
            # ODT, RTF, EPUB
            ("odt","docx"): self._pan,        ("odt","pdf"):  self._doc_via_html_pdf,
            ("odt","txt"):  self._pan,        ("odt","html"): self._pan,
            ("odt","epub"): self._pan,        ("odt","rtf"):  self._pan,
            ("odt","csv"):  self._txt_csv,
            ("rtf","docx"): self._pan,        ("rtf","pdf"):  self._doc_via_html_pdf,
            ("rtf","txt"):  self._pan,        ("rtf","html"): self._pan,
            ("rtf","epub"): self._pan,        ("rtf","odt"):  self._pan,
            ("rtf","csv"):  self._txt_csv,
            ("epub","pdf"): self._doc_via_html_pdf, ("epub","docx"):self._pan,
            ("epub","txt"): self._pan,        ("epub","html"):self._pan,
            ("epub","md"):  self._pan,        ("epub","odt"): self._pan,
            ("epub","rtf"): self._pan,        ("epub","csv"): self._txt_csv,
            # HTML→RTF
            ("html","rtf"): self._pan,
        }
        return table.get((src, tgt), self._pan)

    # ── Helpers ──────────────────────────────────────────────────────────────
    def _write(self, odir, name, text):
        p = odir / name
        p.write_text(text, encoding="utf-8")
        return p

    def _result(self, odir, name):
        p = odir / name
        if not p.exists():
            raise FileNotFoundError(f"No output produced: {p}")
        return {"filename": name, "size": p.stat().st_size}

    def _base_html(self, title, body, extra_css=""):
        esc = _html.escape(title)
        return (f'<!DOCTYPE html><html lang="en"><head><meta charset="utf-8">'
                f'<meta name="viewport" content="width=device-width,initial-scale=1">'
                f'<title>{esc}</title><style>'
                f'body{{font-family:system-ui,sans-serif;max-width:900px;margin:2rem auto;'
                f'line-height:1.7;padding:0 1rem;color:#111}}'
                f'pre,code{{font-family:monospace;background:#f4f4f4;border-radius:4px}}'
                f'pre{{padding:1rem;overflow-x:auto}}code{{padding:.1em .3em}}'
                f'table{{border-collapse:collapse;width:100%;margin:1rem 0}}'
                f'th{{background:#f0f0f0}}th,td{{border:1px solid #ddd;padding:.5rem .75rem}}'
                f'img{{max-width:100%;height:auto;display:block;margin:1rem 0}}'
                f'h1,h2,h3{{letter-spacing:-.02em}}{extra_css}'
                f'</style></head><body>{body}</body></html>')

    # ── MD → HTML ─────────────────────────────────────────────────────────────
    async def _md_html(self, inp, tgt, odir, cb, stem):
        await cb(30, "Rendering Markdown…")
        try:
            import markdown as md_lib
        except ImportError:
            raise RuntimeError("Run: pip install markdown Pygments")
        text = inp.read_text(encoding="utf-8", errors="replace")
        body = md_lib.markdown(text, extensions=["tables","fenced_code","toc","nl2br","attr_list"])
        fn = f"{stem}.html"
        self._write(odir, fn, self._base_html(stem, body))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── MD → PDF ──────────────────────────────────────────────────────────────
    async def _md_pdf(self, inp, tgt, odir, cb, stem):
        await cb(20, "MD → HTML → PDF…")
        try:
            import markdown as md_lib
        except ImportError:
            raise RuntimeError("Run: pip install markdown")
        text = inp.read_text(encoding="utf-8", errors="replace")
        body = md_lib.markdown(text, extensions=["tables","fenced_code","toc","nl2br"])
        html_content = self._base_html(stem, body)
        tmp_html = odir / f"_tmp_{stem}.html"
        tmp_html.write_text(html_content, encoding="utf-8")
        await cb(50, "Rendering to PDF…")
        try:
            r = await self._html_pdf(tmp_html, "pdf", odir, cb, stem)
            return r
        finally:
            tmp_html.unlink(missing_ok=True)

    # ── MD → TXT ──────────────────────────────────────────────────────────────
    async def _md_txt(self, inp, tgt, odir, cb, stem):
        await cb(30, "Stripping Markdown…")
        text = inp.read_text(encoding="utf-8", errors="replace")
        text = re.sub(r'#{1,6}\s', '', text)
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        text = re.sub(r'`(.*?)`', r'\1', text)
        text = re.sub(r'!\[.*?\]\(.*?\)', '', text)
        text = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', text)
        fn = f"{stem}.txt"
        self._write(odir, fn, text)
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── TXT → HTML ────────────────────────────────────────────────────────────
    async def _txt_html(self, inp, tgt, odir, cb, stem):
        await cb(30, "Wrapping text…")
        text = _html.escape(inp.read_text(encoding="utf-8", errors="replace"))
        fn = f"{stem}.html"
        self._write(odir, fn, self._base_html(stem, f"<pre>{text}</pre>"))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── TXT → PDF ─────────────────────────────────────────────────────────────
    async def _txt_pdf(self, inp, tgt, odir, cb, stem):
        await cb(30, "Building PDF…")
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import mm
        except ImportError:
            raise RuntimeError("Run: pip install reportlab")
        fn = f"{stem}.pdf"
        out = odir / fn
        text = inp.read_text(encoding="utf-8", errors="replace")
        doc = SimpleDocTemplate(str(out), pagesize=A4,
                                leftMargin=25*mm, rightMargin=25*mm,
                                topMargin=20*mm, bottomMargin=20*mm)
        styles = getSampleStyleSheet()
        mono = ParagraphStyle("M", parent=styles["Normal"], fontName="Courier", fontSize=9, leading=14)
        story = []
        for line in text.split("\n"):
            s = _sanitize(line)
            story.append(Paragraph(_html.escape(s) if s.strip() else "&nbsp;", mono))
            story.append(Spacer(1,1))
        doc.build(story)
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── TXT → XML ─────────────────────────────────────────────────────────────
    async def _txt_xml(self, inp, tgt, odir, cb, stem):
        await cb(30, "TXT → XML…")
        text = inp.read_text(encoding="utf-8", errors="replace")
        lines = [f"  <line>{_html.escape(_sanitize(l))}</line>" for l in text.split("\n") if l.strip()]
        xml = f'<?xml version="1.0" encoding="UTF-8"?>\n<document>\n' + "\n".join(lines) + "\n</document>"
        fn = f"{stem}.xml"
        self._write(odir, fn, xml)
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── TXT → PPTX ────────────────────────────────────────────────────────────
    async def _txt_pptx(self, inp, tgt, odir, cb, stem):
        await cb(20, "TXT → PPTX…")
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise RuntimeError("Run: pip install python-pptx")
        text = inp.read_text(encoding="utf-8", errors="replace")
        prs = Presentation()
        blank_layout = prs.slide_layouts[1]
        # Split into slides by double newline or every 10 lines
        chunks = re.split(r'\n{2,}', text.strip())
        if not chunks: chunks = [text]
        for i, chunk in enumerate(chunks[:50]):  # max 50 slides
            lines = [l.strip() for l in chunk.split("\n") if l.strip()]
            if not lines: continue
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.title.text = lines[0][:80]
            if len(lines) > 1 and slide.placeholders[1]:
                slide.placeholders[1].text = "\n".join(lines[1:10])
            await cb(20 + int(i/len(chunks)*70), f"Slide {i+1}…")
        fn = f"{stem}.pptx"
        prs.save(str(odir / fn))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── HTML → PDF ────────────────────────────────────────────────────────────
    async def _html_pdf(self, inp, tgt, odir, cb, stem):
        await cb(30, "Rendering HTML → PDF (with images)…")
        try:
            from weasyprint import HTML, CSS
            fn = f"{stem}.pdf"
            out = odir / fn
            HTML(filename=str(inp)).write_pdf(str(out),
                stylesheets=[CSS(string="@page{margin:20mm} img{max-width:100%}")])
            await cb(95, "Done!")
            return self._result(odir, fn)
        except ImportError:
            return await self._html_pdf_rl(inp, tgt, odir, cb, stem)
        except Exception as e:
            logger.warning(f"WeasyPrint failed: {e}, fallback to reportlab")
            return await self._html_pdf_rl(inp, tgt, odir, cb, stem)

    async def _html_pdf_rl(self, inp, tgt, odir, cb, stem):
        await cb(40, "HTML→PDF fallback (text extraction)…")
        from html.parser import HTMLParser
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import mm

        class TX(HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts = []; self._skip = False
            def handle_starttag(self, t, a):
                if t in ('script','style'): self._skip = True
            def handle_endtag(self, t):
                if t in ('script','style'): self._skip = False
            def handle_data(self, d):
                if not self._skip and d.strip(): self.parts.append(d.strip())
        p = TX(); p.feed(inp.read_text(encoding="utf-8", errors="replace"))
        text_content = "\n".join(p.parts)
        fn = f"{stem}.pdf"
        out = odir / fn
        doc = SimpleDocTemplate(str(out), pagesize=A4,
                                leftMargin=25*mm, rightMargin=25*mm,
                                topMargin=20*mm, bottomMargin=20*mm)
        styles = getSampleStyleSheet()
        story = [Paragraph(_html.escape(_sanitize(l)) if l.strip() else "&nbsp;", styles["Normal"])
                 for l in text_content.split("\n")]
        doc.build(story)
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── HTML → TXT ────────────────────────────────────────────────────────────
    async def _html_txt(self, inp, tgt, odir, cb, stem):
        await cb(30, "Extracting text from HTML…")
        from html.parser import HTMLParser
        class TX(HTMLParser):
            def __init__(self): super().__init__(); self.parts=[]; self._skip=False
            def handle_starttag(self,t,a):
                if t in('script','style'): self._skip=True
                if t in('br','p','div','h1','h2','h3','h4','li'): self.parts.append('\n')
            def handle_endtag(self,t):
                if t in('script','style'): self._skip=False
            def handle_data(self,d):
                if not self._skip and d.strip(): self.parts.append(d.strip())
        p = TX(); p.feed(inp.read_text(encoding="utf-8", errors="replace"))
        fn = f"{stem}.txt"
        self._write(odir, fn, "\n".join(p.parts))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── HTML → XML ────────────────────────────────────────────────────────────
    async def _html_xml(self, inp, tgt, odir, cb, stem):
        await cb(30, "HTML → XML…")
        text = inp.read_text(encoding="utf-8", errors="replace")
        fn = f"{stem}.xml"
        self._write(odir, fn, f'<?xml version="1.0" encoding="UTF-8"?>\n<html>{text}</html>')
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── DOCX → TXT ────────────────────────────────────────────────────────────
    async def _docx_txt(self, inp, tgt, odir, cb, stem):
        await cb(30, "Extracting DOCX text…")
        try:
            from docx import Document
        except ImportError:
            raise RuntimeError("Run: pip install python-docx")
        doc = Document(str(inp))
        parts = [p.text for p in doc.paragraphs]
        for tbl in doc.tables:
            for row in tbl.rows:
                parts.append("\t".join(c.text for c in row.cells))
        fn = f"{stem}.txt"
        self._write(odir, fn, "\n".join(parts))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── DOCX → HTML (with embedded images) ───────────────────────────────────
    async def _docx_html(self, inp, tgt, odir, cb, stem):
        await cb(30, "DOCX → HTML (with images)…")
        try:
            import mammoth
            # Custom image handler — embed as base64
            def convert_image(image):
                with image.open() as img_file:
                    img_bytes = img_file.read()
                mime = image.content_type or "image/png"
                b64 = base64.b64encode(img_bytes).decode()
                return {"src": f"data:{mime};base64,{b64}"}

            fn = f"{stem}.html"
            with open(inp, "rb") as f:
                result = mammoth.convert_to_html(f, convert_image=mammoth.images.img_element(convert_image))
            body = result.value
            self._write(odir, fn, self._base_html(stem, body))
            await cb(95, "Done!")
            return self._result(odir, fn)
        except ImportError:
            return await self._pan(inp, tgt, odir, cb, stem)

    # ── DOCX → PDF (with images via WeasyPrint) ───────────────────────────────
    async def _docx_pdf(self, inp, tgt, odir, cb, stem):
        """DOCX → PDF: LibreOffice headless (best fidelity) → WeasyPrint → Pandoc fallback."""
        await cb(20, "DOCX → PDF…")
        fn = f"{stem}.pdf"
        out = odir / fn

        # Strategy 1: LibreOffice headless — preserves layout perfectly
        lo = shutil.which("libreoffice") or shutil.which("soffice")
        if lo:
            try:
                await cb(30, "Converting with LibreOffice…")
                proc = await asyncio.create_subprocess_exec(
                    lo, "--headless", "--norestore", "--nofirststartwizard",
                    "--convert-to", "pdf",
                    "--outdir", str(odir), str(inp),
                    stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
                )
                _, err = await asyncio.wait_for(proc.communicate(), timeout=15)
                if proc.returncode == 0:
                    lo_out = odir / f"{inp.stem}.pdf"
                    if lo_out.exists() and lo_out != out:
                        lo_out.rename(out)
                    if out.exists():
                        await cb(95, "Done!")
                        return self._result(odir, fn)
                else:
                    logger.warning(f"LibreOffice failed: {err.decode(errors='replace')[-200:]}")
            except Exception as e:
                logger.warning(f"LibreOffice DOCX→PDF failed ({e}), trying WeasyPrint")

        # Strategy 2: mammoth + WeasyPrint
        try:
            import mammoth
            from weasyprint import HTML, CSS

            def convert_image(image):
                with image.open() as img_file:
                    img_bytes = img_file.read()
                mime = image.content_type or "image/png"
                b64 = base64.b64encode(img_bytes).decode()
                return {"src": f"data:{mime};base64,{b64}"}

            await cb(40, "Converting with WeasyPrint…")
            with open(inp, "rb") as f:
                result = mammoth.convert_to_html(f, convert_image=mammoth.images.img_element(convert_image))
            body = result.value
            html_content = self._base_html(stem, body)
            HTML(string=html_content).write_pdf(str(out),
                stylesheets=[CSS(string=(
                    "@page{margin:20mm}"
                    " img{max-width:100%;height:auto}"
                    " table{width:100%;border-collapse:collapse}"
                    " td,th{border:1px solid #ddd;padding:4px}"
                ))])
            if out.exists():
                await cb(95, "Done!")
                return self._result(odir, fn)
        except ImportError as e:
            logger.warning(f"mammoth/weasyprint not available ({e}), trying Pandoc")
        except Exception as e:
            logger.warning(f"WeasyPrint DOCX→PDF failed ({e}), trying Pandoc")

        # Strategy 3: Pandoc fallback
        return await self._pan(inp, tgt, odir, cb, stem)

    # ── DOCX → PPTX ───────────────────────────────────────────────────────────
    async def _docx_pptx(self, inp, tgt, odir, cb, stem):
        await cb(20, "DOCX → PPTX…")
        try:
            from docx import Document
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise RuntimeError("Run: pip install python-docx python-pptx")
        doc = Document(str(inp))
        prs = Presentation()
        layout = prs.slide_layouts[1]
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        # Group paragraphs into slides (1 heading + ~5 bullets per slide)
        slides = []
        current = []
        for p in paras:
            current.append(p)
            if len(current) >= 6:
                slides.append(current); current = []
        if current: slides.append(current)
        for i, group in enumerate(slides[:50]):
            await cb(20 + int(i/len(slides)*70), f"Slide {i+1}…")
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = group[0][:80]
            if len(group) > 1 and slide.placeholders[1]:
                slide.placeholders[1].text = "\n".join(group[1:])
        fn = f"{stem}.pptx"
        prs.save(str(odir / fn))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── PDF → TXT (Unicode-safe) ───────────────────────────────────────────────
    async def _pdf_txt(self, inp, tgt, odir, cb, stem):
        """Extract text from PDF — handles Unicode, Hindi/Devanagari, Arabic, CJK via pypdfium2."""
        await cb(20, "Extracting PDF text (Unicode-safe)…")
        fn = f"{stem}.txt"

        # Strategy 1: pypdfium2 — best Unicode support (correct font encoding)
        try:
            import pypdfium2 as pdfium
            pdf = pdfium.PdfDocument(str(inp))
            total = len(pdf)
            if total == 0:
                raise RuntimeError("PDF has no pages or is password-protected.")
            pages = []
            # Process pages in batches of 10 — reduces await overhead
            for i in range(total):
                page = pdf[i]
                textpage = page.get_textpage()
                text = textpage.get_text_range() or ""
                pages.append(_sanitize(text))
                textpage.close()
                page.close()
                # Only update progress every 10 pages to reduce WS overhead
                if i == 0 or i % 10 == 0 or i == total - 1:
                    await cb(20 + int(i / total * 70), f"Page {i+1}/{total}…")
            pdf.close()
            content = "\n\n---\n\n".join(pages)
            if len(content.strip()) > 20:
                self._write(odir, fn, content)
                await cb(95, "Done!")
                return self._result(odir, fn)
        except ImportError:
            logger.debug("pypdfium2 not available, using pypdf fallback")
        except Exception as e:
            logger.warning(f"pypdfium2 text extraction failed ({e}), falling back")

        # Strategy 2: pypdf fallback
        from pypdf import PdfReader
        reader = PdfReader(str(inp))
        total = len(reader.pages)
        if total == 0:
            raise RuntimeError("PDF has no pages or is password-protected.")
        pages = []
        for i, page in enumerate(reader.pages):
            raw = page.extract_text() or ""
            pages.append(_sanitize(raw))
            await cb(20 + int(i/total*70), f"Page {i+1}/{total}…")
        self._write(odir, fn, "\n\n---\n\n".join(pages))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── PDF → DOCX via page rendering (pypdfium2) ───────────────────────────
    async def _pdf_docx_render(self, inp, tgt, odir, cb, stem):
        """Render each PDF page as 150 DPI image and embed in DOCX.
        Uses pypdfium2 — renders EXACTLY what the PDF looks like.
        Zero broken images guaranteed."""
        import pypdfium2 as pdfium
        from pypdf import PdfReader
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches, Cm
        from PIL import Image as PILImage

        pdf    = pdfium.PdfDocument(str(inp))
        reader = PdfReader(str(inp))
        total  = len(pdf)

        doc = Document()
        for section in doc.sections:
            section.top_margin = Cm(1.5); section.bottom_margin = Cm(1.5)
            section.left_margin = Cm(2.0); section.right_margin = Cm(2.0)
        doc.styles["Normal"].font.name = "Calibri"
        doc.styles["Normal"].font.size = Pt(10)

        for i in range(total):
            await cb(20 + int(i / total * 70), f"Rendering page {i+1}/{total}…")
            if i > 0:
                doc.add_page_break()

            h = doc.add_heading(f"Page {i+1}", level=2)
            if h.runs:
                h.runs[0].font.color.rgb = RGBColor(0x5B, 0x5F, 0xEC)

            # Render full PDF page at 150 DPI — text, images, graphics all included
            try:
                page = pdf[i]
                bitmap = page.render(scale=150/72, rotation=0)
                pil = bitmap.to_pil()
                if pil.mode != "RGB":
                    pil = pil.convert("RGB")
                img_buf = io.BytesIO()
                pil.save(img_buf, "JPEG", quality=85, optimize=True)
                img_buf.seek(0)
                page_w_in = min(6.7, pil.width / 150.0)
                doc.add_picture(img_buf, width=Inches(page_w_in))
            except Exception as e:
                logger.warning(f"Page {i+1} render error: {e}")
                doc.add_paragraph(f"[Page {i+1} could not be rendered]")

            # Add extracted text below the image for searchability
            try:
                raw = reader.pages[i].extract_text() or ""
                text = _sanitize(raw)
                if text.strip():
                    lbl = doc.add_paragraph()
                    run = lbl.add_run("Extracted text:")
                    run.bold = True; run.font.color.rgb = RGBColor(0x88,0x88,0x88)
                    for line in text.split("\n"):
                        line = line.strip()
                        if line:
                            try: doc.add_paragraph(line).style = doc.styles["Normal"]
                            except Exception: pass
            except Exception:
                pass

        pdf.close()
        fn = f"{stem}.docx"
        doc.save(str(odir / fn))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── PDF → HTML (with embedded images) ────────────────────────────────────
    async def _pdf_html(self, inp, tgt, odir, cb, stem):
        await cb(20, "PDF → HTML (rendering pages)…")
        try:
            import pypdfium2 as pdfium
            from pypdf import PdfReader
            pdf    = pdfium.PdfDocument(str(inp))
            reader = PdfReader(str(inp))
            total  = len(pdf)
            sections = []
            for i in range(total):
                await cb(20 + int(i/total*60), f"Page {i+1}/{total}…")
                try:
                    page   = pdf[i]
                    bitmap = page.render(scale=120/72, rotation=0)
                    pil    = bitmap.to_pil().convert("RGB")
                    buf    = io.BytesIO()
                    pil.save(buf, "JPEG", quality=80, optimize=True)
                    b64    = base64.b64encode(buf.getvalue()).decode()
                    img_tag = f'<img src="data:image/jpeg;base64,{b64}" alt="page {i+1}" style="max-width:100%;display:block;margin:0 auto 1rem;border:1px solid #eee">'
                except Exception:
                    img_tag = ""
                text = _sanitize(reader.pages[i].extract_text() or "")
                txt_html = f"<details><summary>Text</summary><pre style='white-space:pre-wrap;font-size:.8rem;color:#666'>{_html.escape(text)}</pre></details>" if text.strip() else ""
                sections.append(f'<div class="page"><h3>Page {i+1}</h3>{img_tag}{txt_html}</div>')
            pdf.close()
            fn = f"{stem}.html"
            self._write(odir, fn, self._base_html(stem, "\n".join(sections),
                ".page{border-bottom:2px solid #eee;padding-bottom:2rem;margin-bottom:2rem}"))
            await cb(95, "Done!")
            return self._result(odir, fn)
        except ImportError:
            pass
        except Exception as e:
            logger.warning(f"pypdfium2 HTML render failed: {e}")
        # Fallback: pdfminer text
        try:
            from pdfminer.high_level import extract_text
            text = extract_text(str(inp))
        except ImportError:
            from pypdf import PdfReader
            text = "\n".join(p.extract_text() or "" for p in PdfReader(str(inp)).pages)
        fn = f"{stem}.html"
        self._write(odir, fn, self._base_html(stem, f"<pre>{_html.escape(_sanitize(text))}</pre>"))
        await cb(95, "Done!")
        return self._result(odir, fn)


    # ── PDF → DOCX (with images + Unicode sanitize) ───────────────────────────
    async def _pdf_docx(self, inp, tgt, odir, cb, stem):
        await cb(20, "PDF → DOCX (rendering pages with images)…")
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches
        except ImportError as e:
            raise RuntimeError(f"Run: pip install python-docx\nError: {e}")

        # Try pypdfium2 first — renders full page as image (perfect quality, no broken images)
        try:
            return await self._pdf_docx_render(inp, tgt, odir, cb, stem)
        except ImportError:
            pass  # fall through to pypdf extraction approach
        except Exception as e:
            logger.warning(f"pypdfium2 render failed ({e}), falling back to extraction")

        def _get_content_images(page):
            """
            Extract only real content images from a PDF page.
            Skips: full-page backgrounds, tiny decorations, duplicates.
            Returns list of (BytesIO_of_PNG, width_in_inches).
            """
            if not hasattr(page, 'images'):
                return []
            mb   = page.mediabox
            pg_w = float(mb.width)
            pg_h = float(mb.height)
            # Rough pixel dimensions at 150 DPI
            est_px_w = pg_w * (150 / 72)
            est_px_h = pg_h * (150 / 72)

            seen  = set()
            found = []

            for img_obj in page.images:
                try:
                    data = img_obj.data
                    if not data or len(data) < 200:
                        continue
                    # Deduplicate by content hash
                    sig = hash(data[:512])
                    if sig in seen:
                        continue
                    seen.add(sig)

                    pil = PILImage.open(io.BytesIO(data))
                    pil.load()
                    w, h = pil.size

                    # Skip tiny images (bullets, icons, watermarks)
                    if w < 40 or h < 40:
                        continue

                    # Skip full-page background images:
                    # An image is a background if it fills >75% of the page in BOTH axes
                    if w > est_px_w * 0.75 and h > est_px_h * 0.75:
                        logger.debug(f"Skip background image {w}x{h} (page ~{est_px_w:.0f}x{est_px_h:.0f})")
                        continue

                    # Convert to RGB PNG (handles LA, RGBA, P, CMYK, JP2 etc.)
                    out = io.BytesIO()
                    if pil.mode == 'CMYK':
                        pil = pil.convert('RGB')
                    elif pil.mode in ('RGBA', 'LA', 'PA'):
                        bg = PILImage.new('RGB', pil.size, (255, 255, 255))
                        if pil.mode == 'PA':
                            pil = pil.convert('RGBA')
                        elif pil.mode == 'LA':
                            pil = pil.convert('RGBA')
                        bg.paste(pil, mask=pil.split()[-1])
                        pil = bg
                    elif pil.mode != 'RGB':
                        pil = pil.convert('RGB')

                    pil.save(out, 'PNG', optimize=True)
                    out.seek(0)

                    # Width in inches for DOCX — use actual pixel size at 150 DPI
                    # Cap at 5.5" (A4 with 2.5cm margins each side)
                    w_inches = min(5.5, w / 150.0)
                    # Minimum useful size: 0.5 inches
                    if w_inches < 0.5:
                        continue

                    found.append((out, w_inches, w * h))

                except Exception as e:
                    logger.debug(f"Image decode skipped: {e}")
                    continue

            # Sort by pixel area (largest first = most significant content)
            found.sort(key=lambda x: x[2], reverse=True)
            # Return up to 6 content images per page
            return [(buf, w_in) for buf, w_in, _ in found[:6]]

        from pypdf import PdfReader  # noqa — needed in fallback path
        reader = PdfReader(str(inp))
        total = len(reader.pages)
        if total == 0:
            raise RuntimeError("PDF has no pages or is password-protected.")

        doc = Document()
        doc.styles["Normal"].font.name = "Calibri"
        doc.styles["Normal"].font.size = Pt(11)

        for i, page in enumerate(reader.pages):
            await cb(20 + int(i / total * 70), f"Page {i+1}/{total}…")
            text = _sanitize(page.extract_text() or "")

            if i > 0:
                doc.add_page_break()

            # Page heading
            h = doc.add_heading(f"Page {i+1}", level=2)
            if h.runs:
                h.runs[0].font.color.rgb = RGBColor(0x5B, 0x5F, 0xEC)

            # ── Insert content images ─────────────────────────────────────
            images = _get_content_images(page)
            logger.info(f"Page {i+1}: {len(images)} content images found")
            for img_buf, w_in in images:
                try:
                    doc.add_picture(img_buf, width=Inches(w_in))
                except Exception as e:
                    logger.debug(f"add_picture failed: {e}")

            # ── Insert text ───────────────────────────────────────────────
            for line in text.split("\n"):
                line = line.strip()
                if not line:
                    continue
                # Detect headings: short, all-caps, meaningful length
                if 4 < len(line) < 80 and line.isupper():
                    try:
                        doc.add_heading(line.title(), level=3)
                        continue
                    except Exception:
                        pass
                try:
                    p = doc.add_paragraph(line)
                    p.style = doc.styles["Normal"]
                except Exception:
                    pass

        fn = f"{stem}.docx"
        doc.save(str(odir / fn))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── PDF → MD ──────────────────────────────────────────────────────────────
    async def _pdf_md(self, inp, tgt, odir, cb, stem):
        """PDF → Markdown: extract text then format as clean Markdown."""
        await cb(20, "PDF → Markdown…")
        fn = f"{stem}.md"
        txt_result = await self._pdf_txt(inp, "txt", odir, cb, stem)
        txt_path = odir / txt_result["filename"]
        raw = txt_path.read_text(encoding="utf-8", errors="replace")
        txt_path.unlink(missing_ok=True)
        lines = raw.split("\n")
        md_lines = []
        for line in lines:
            stripped = line.strip()
            if not stripped:
                md_lines.append("")
                continue
            if stripped == "---":
                md_lines.append("\n---\n")
                continue
            if 4 < len(stripped) < 80 and stripped.isupper():
                md_lines.append(f"## {stripped.title()}")
            else:
                md_lines.append(stripped)
        self._write(odir, fn, "\n".join(md_lines))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── PDF → XML ──────────────────────────────────────────────────────────────
    async def _pdf_xml(self, inp, tgt, odir, cb, stem):
        await cb(20, "PDF → XML…")
        txt_result = await self._pdf_txt(inp, "txt", odir, cb, stem)
        txt_path = odir / txt_result["filename"]
        raw = txt_path.read_text(encoding="utf-8", errors="replace")
        txt_path.unlink(missing_ok=True)
        pages = raw.split("\n\n---\n\n")
        pages_xml = [f'  <page number="{i+1}">{_html.escape(p)}</page>' for i,p in enumerate(pages)]
        xml = '<?xml version="1.0" encoding="UTF-8"?>\n<document>\n' + "\n".join(pages_xml) + "\n</document>"
        fn = f"{stem}.xml"
        self._write(odir, fn, xml)
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── PPTX → DOCX (via text extraction) ───────────────────────────────────
    async def _pptx_docx(self, inp, tgt, odir, cb, stem):
        """PPTX → DOCX: extract text per slide and build a Word document."""
        await cb(20, "PPTX → DOCX…")
        try:
            from pptx import Presentation
            from docx import Document
            from docx.shared import Pt, RGBColor
        except ImportError:
            raise RuntimeError("Run: pip install python-pptx python-docx")
        prs = Presentation(str(inp))
        doc = Document()
        doc.styles["Normal"].font.name = "Calibri"
        doc.styles["Normal"].font.size = Pt(11)
        total = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            await cb(20 + int(i/total*70), f"Slide {i+1}/{total}…")
            h = doc.add_heading(f"Slide {i+1}", level=1)
            if h.runs:
                h.runs[0].font.color.rgb = RGBColor(0x5B, 0x5F, 0xEC)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = _sanitize(para.text.strip())
                        if t:
                            try:
                                p = doc.add_paragraph(t)
                                p.style = doc.styles["Normal"]
                            except Exception:
                                pass
        fn = f"{stem}.docx"
        doc.save(str(odir / fn))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── PPTX → TXT ────────────────────────────────────────────────────────────
    async def _pptx_txt(self, inp, tgt, odir, cb, stem):
        await cb(30, "Extracting PPTX text…")
        try:
            from pptx import Presentation
        except ImportError:
            raise RuntimeError("Run: pip install python-pptx")
        prs = Presentation(str(inp))
        lines = []
        for i, slide in enumerate(prs.slides):
            lines.append(f"\n=== Slide {i+1} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = _sanitize(para.text.strip())
                        if t: lines.append(t)
        fn = f"{stem}.txt"
        self._write(odir, fn, "\n".join(lines))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── PPTX → HTML (with embedded images) ───────────────────────────────────
    async def _pptx_html(self, inp, tgt, odir, cb, stem):
        await cb(30, "PPTX → HTML (with images)…")
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise RuntimeError("Run: pip install python-pptx")
        prs = Presentation(str(inp))
        slides_html = []
        total = len(prs.slides)
        for i, slide in enumerate(prs.slides):
            await cb(30 + int(i/total*60), f"Slide {i+1}/{total}…")
            parts = [f'<h3>Slide {i+1}</h3>']
            # Extract text
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = _sanitize(para.text.strip())
                        if t: parts.append(f"<p>{_html.escape(t)}</p>")
                # Extract images
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_bytes = shape.image.blob
                        mime = shape.image.content_type or "image/png"
                        b64 = base64.b64encode(img_bytes).decode()
                        parts.append(f'<img src="data:{mime};base64,{b64}" alt="slide {i+1}" style="max-width:100%;margin:.5rem 0">')
                except Exception:
                    pass
            slides_html.append(
                f'<section style="border:1px solid #ddd;padding:1.5rem;margin-bottom:1.5rem;border-radius:8px">'
                + "".join(parts) + "</section>"
            )
        fn = f"{stem}.html"
        self._write(odir, fn, self._base_html(stem, "".join(slides_html)))
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── PPTX → PDF (with images) ──────────────────────────────────────────────
    async def _pptx_pdf(self, inp, tgt, odir, cb, stem):
        await cb(20, "PPTX → PDF (with images)…")
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Image as RLImage
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import mm
            from reportlab.lib.colors import HexColor
        except ImportError:
            raise RuntimeError("Run: pip install python-pptx reportlab")

        prs = Presentation(str(inp))
        fn = f"{stem}.pdf"
        out = odir / fn
        doc = SimpleDocTemplate(str(out), pagesize=landscape(A4),
                                leftMargin=20*mm, rightMargin=20*mm,
                                topMargin=15*mm, bottomMargin=15*mm)
        styles = getSampleStyleSheet()
        title_s = ParagraphStyle("ST", parent=styles["Heading1"], fontSize=18,
                                  textColor=HexColor("#2d2d5e"), spaceAfter=10)
        body_s  = ParagraphStyle("SB", parent=styles["Normal"], fontSize=11, leading=16)
        num_s   = ParagraphStyle("SN", parent=styles["Normal"], fontSize=8,
                                  textColor=HexColor("#999999"))
        story = []
        total = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            await cb(20 + int(i/total*70), f"Slide {i+1}/{total}…")
            story.append(Paragraph(f"Slide {i+1}", num_s))
            story.append(Spacer(1, 4))
            title_text = ""
            body_items = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for j, para in enumerate(shape.text_frame.paragraphs):
                        t = _sanitize(para.text.strip())
                        if not t: continue
                        if j == 0 and not title_text:
                            title_text = t
                        else:
                            body_items.append(t)
                # Images
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img_buf = io.BytesIO(shape.image.blob)
                        from PIL import Image as PILImage
                        pil = PILImage.open(img_buf)
                        # Max 4 inches wide in landscape
                        w_in = min(4, pil.width / 96)
                        h_in = w_in * pil.height / pil.width
                        out_buf = io.BytesIO()
                        pil.convert("RGB").save(out_buf, "PNG"); out_buf.seek(0)
                        story.append(RLImage(out_buf, width=w_in*72, height=h_in*72))
                except Exception:
                    pass

            if title_text:
                story.append(Paragraph(_html.escape(_sanitize(title_text)), title_s))
            for bt in body_items:
                story.append(Paragraph(f"• {_html.escape(_sanitize(bt))}", body_s))
            if i < total - 1:
                story.append(PageBreak())

        doc.build(story)
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── XLSX conversions ──────────────────────────────────────────────────────
    async def _xlsx_csv(self, inp, tgt, odir, cb, stem):
        await cb(30, "XLSX → CSV…")
        try: import openpyxl
        except ImportError: raise RuntimeError("Run: pip install openpyxl")
        wb = openpyxl.load_workbook(str(inp), read_only=True, data_only=True)
        fn = f"{stem}.csv"
        with open(odir/fn, "w", newline="", encoding="utf-8") as f:
            w = _csv.writer(f)
            for row in wb.active.iter_rows(values_only=True):
                w.writerow([str(c) if c is not None else "" for c in row])
        wb.close(); await cb(95, "Done!")
        return self._result(odir, fn)

    async def _xlsx_txt(self, inp, tgt, odir, cb, stem):
        await cb(30, "XLSX → TXT…")
        try: import openpyxl
        except ImportError: raise RuntimeError("Run: pip install openpyxl")
        wb = openpyxl.load_workbook(str(inp), read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"\n=== {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                lines.append("\t".join(str(c) if c is not None else "" for c in row))
        wb.close()
        fn = f"{stem}.txt"
        self._write(odir, fn, "\n".join(lines))
        await cb(95, "Done!")
        return self._result(odir, fn)

    async def _xlsx_html(self, inp, tgt, odir, cb, stem):
        await cb(30, "XLSX → HTML…")
        try: import openpyxl
        except ImportError: raise RuntimeError("Run: pip install openpyxl")
        wb = openpyxl.load_workbook(str(inp), read_only=True, data_only=True)
        tables = []
        for sheet in wb.worksheets:
            rows_html = []
            first = True
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                tag = "th" if first else "td"
                row_html = "".join(f"<{tag}>{_html.escape(c)}</{tag}>" for c in cells)
                rows_html.append(f"<tr>{row_html}</tr>"); first = False
            tables.append(f'<h2>{_html.escape(sheet.title)}</h2><table>{"".join(rows_html)}</table>')
        wb.close()
        fn = f"{stem}.html"
        self._write(odir, fn, self._base_html(stem, "".join(tables)))
        await cb(95, "Done!")
        return self._result(odir, fn)

    async def _xlsx_pdf(self, inp, tgt, odir, cb, stem):
        await cb(20, "XLSX → PDF…")
        try:
            import openpyxl
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import mm
            from reportlab.lib import colors
        except ImportError:
            raise RuntimeError("Run: pip install openpyxl reportlab")
        wb = openpyxl.load_workbook(str(inp), read_only=True, data_only=True)
        fn = f"{stem}.pdf"
        out = odir / fn
        doc = SimpleDocTemplate(str(out), pagesize=landscape(A4),
                                leftMargin=15*mm, rightMargin=15*mm,
                                topMargin=15*mm, bottomMargin=15*mm)
        styles = getSampleStyleSheet()
        story = []
        for si, sheet in enumerate(wb.worksheets):
            await cb(20 + int(si/len(wb.worksheets)*70), f"Sheet: {sheet.title}…")
            story.append(Paragraph(f"Sheet: {_html.escape(sheet.title)}", styles["Heading2"]))
            story.append(Spacer(1, 6))
            data = [[str(c)[:40] if c is not None else "" for c in row]
                    for row in sheet.iter_rows(values_only=True)]
            if data:
                t = Table(data, repeatRows=1)
                t.setStyle(TableStyle([
                    ("BACKGROUND",(0,0),(-1,0),colors.HexColor("#e8e8f0")),
                    ("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),
                    ("FONTSIZE",(0,0),(-1,-1),8),
                    ("GRID",(0,0),(-1,-1),0.5,colors.HexColor("#cccccc")),
                    ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f8f8fc")]),
                    ("TOPPADDING",(0,0),(-1,-1),3),("BOTTOMPADDING",(0,0),(-1,-1),3),
                ]))
                story.append(t)
            story.append(Spacer(1, 20))
        wb.close()
        doc.build(story)
        await cb(95, "Done!")
        return self._result(odir, fn)

    async def _xlsx_json(self, inp, tgt, odir, cb, stem):
        await cb(30, "XLSX → JSON…")
        import json
        try: import openpyxl
        except ImportError: raise RuntimeError("Run: pip install openpyxl")
        wb = openpyxl.load_workbook(str(inp), read_only=True, data_only=True)
        result = {}
        for sheet in wb.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            if not rows: result[sheet.title]=[]; continue
            headers = [str(h) if h is not None else f"col_{i}" for i,h in enumerate(rows[0])]
            result[sheet.title] = [{headers[i]:(v if v is not None else "") for i,v in enumerate(row)} for row in rows[1:]]
        wb.close()
        fn = f"{stem}.json"
        (odir/fn).write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
        await cb(95, "Done!")
        return self._result(odir, fn)

    async def _xlsx_xml(self, inp, tgt, odir, cb, stem):
        await cb(30, "XLSX → XML…")
        try: import openpyxl
        except ImportError: raise RuntimeError("Run: pip install openpyxl")
        wb = openpyxl.load_workbook(str(inp), read_only=True, data_only=True)
        sheets_xml = []
        for sheet in wb.worksheets:
            rows_xml = []
            rows = list(sheet.iter_rows(values_only=True))
            headers = [str(h) if h is not None else f"col{i}" for i,h in enumerate(rows[0])] if rows else []
            for row in rows[1:]:
                cells = "".join(f"<{_safe_tag(headers[i])}>{_html.escape(str(v) if v else '')}</{_safe_tag(headers[i])}>" for i,v in enumerate(row))
                rows_xml.append(f"  <row>{cells}</row>")
            sheets_xml.append(f'<sheet name="{_html.escape(sheet.title)}">\n' + "\n".join(rows_xml) + "\n</sheet>")
        wb.close()
        xml = '<?xml version="1.0" encoding="UTF-8"?>\n<workbook>\n' + "\n".join(sheets_xml) + "\n</workbook>"
        fn = f"{stem}.xml"
        self._write(odir, fn, xml)
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ── CSV conversions ───────────────────────────────────────────────────────
    async def _csv_xlsx(self, inp, tgt, odir, cb, stem):
        await cb(30, "CSV → XLSX…")
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill
        except ImportError:
            raise RuntimeError("Run: pip install openpyxl")
        wb = openpyxl.Workbook(); ws = wb.active; ws.title = stem[:31]
        with open(inp,"r",encoding="utf-8-sig",errors="replace") as f:
            for row in _csv.reader(f): ws.append(row)
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="E8E8F0")
        fn = f"{stem}.xlsx"; wb.save(str(odir/fn)); await cb(95,"Done!")
        return self._result(odir, fn)

    async def _csv_txt(self, inp, tgt, odir, cb, stem):
        await cb(30, "CSV → TXT…")
        fn = f"{stem}.txt"
        self._write(odir, fn, inp.read_text(encoding="utf-8-sig", errors="replace"))
        await cb(95, "Done!")
        return self._result(odir, fn)

    async def _csv_html(self, inp, tgt, odir, cb, stem):
        await cb(30, "CSV → HTML…")
        rows = []
        with open(inp,"r",encoding="utf-8-sig",errors="replace") as f:
            for row in _csv.reader(f): rows.append(row)
        if not rows: body = "<p>Empty CSV</p>"
        else:
            header = "".join(f"<th>{_html.escape(c)}</th>" for c in rows[0])
            body_rows = "".join(
                "<tr>"+"".join(f"<td>{_html.escape(c)}</td>" for c in row)+"</tr>"
                for row in rows[1:])
            body = f"<table><thead><tr>{header}</tr></thead><tbody>{body_rows}</tbody></table>"
        fn = f"{stem}.html"
        self._write(odir, fn, self._base_html(stem, body)); await cb(95,"Done!")
        return self._result(odir, fn)

    async def _csv_json(self, inp, tgt, odir, cb, stem):
        await cb(30, "CSV → JSON…")
        import json
        rows = []
        with open(inp,"r",encoding="utf-8-sig",errors="replace") as f:
            for row in _csv.DictReader(f): rows.append(dict(row))
        fn = f"{stem}.json"
        (odir/fn).write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")
        await cb(95,"Done!")
        return self._result(odir, fn)

    async def _csv_pdf(self, inp, tgt, odir, cb, stem):
        await cb(20, "CSV → PDF…")
        try:
            from reportlab.lib.pagesizes import A4, landscape
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
            from reportlab.lib.units import mm
            from reportlab.lib import colors
        except ImportError:
            raise RuntimeError("Run: pip install reportlab")
        rows = []
        with open(inp,"r",encoding="utf-8-sig",errors="replace") as f:
            for row in _csv.reader(f):
                rows.append([v[:40]+"…" if len(v)>40 else v for v in row])
        fn = f"{stem}.pdf"; out = odir/fn
        doc = SimpleDocTemplate(str(out), pagesize=landscape(A4),
                                leftMargin=15*mm, rightMargin=15*mm,
                                topMargin=15*mm, bottomMargin=15*mm)
        if rows:
            t = Table(rows, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND",(0,0),(-1,0),colors.HexColor("#e8e8f0")),
                ("FONTNAME",(0,0),(-1,0),"Helvetica-Bold"),
                ("FONTSIZE",(0,0),(-1,-1),8),
                ("GRID",(0,0),(-1,-1),0.5,colors.HexColor("#cccccc")),
                ("ROWBACKGROUNDS",(0,1),(-1,-1),[colors.white,colors.HexColor("#f8f8fc")]),
                ("TOPPADDING",(0,0),(-1,-1),3),("BOTTOMPADDING",(0,0),(-1,-1),3),
            ]))
            doc.build([t])
        await cb(95,"Done!")
        return self._result(odir, fn)

    async def _csv_xml(self, inp, tgt, odir, cb, stem):
        await cb(30, "CSV → XML…")
        rows = []
        with open(inp,"r",encoding="utf-8-sig",errors="replace") as f:
            for row in _csv.DictReader(f): rows.append(dict(row))
        items = []
        for row in rows:
            cells = "".join(f"<{_safe_tag(k)}>{_html.escape(str(v))}</{_safe_tag(k)}>" for k,v in row.items())
            items.append(f"  <record>{cells}</record>")
        xml = '<?xml version="1.0" encoding="UTF-8"?>\n<records>\n' + "\n".join(items) + "\n</records>"
        fn = f"{stem}.xml"
        self._write(odir, fn, xml); await cb(95,"Done!")
        return self._result(odir, fn)

    # ── JSON conversions ──────────────────────────────────────────────────────
    async def _json_csv(self, inp, tgt, odir, cb, stem):
        await cb(30, "JSON → CSV…")
        import json
        with open(inp,encoding="utf-8",errors="replace") as f: data = json.load(f)
        fn = f"{stem}.csv"; out = odir/fn
        if isinstance(data,list) and data and isinstance(data[0],dict):
            headers = list(data[0].keys())
            with open(out,"w",newline="",encoding="utf-8") as f:
                w = _csv.writer(f); w.writerow(headers)
                for row in data: w.writerow([str(row.get(h,"")) for h in headers])
        else:
            with open(out,"w",newline="",encoding="utf-8") as f:
                w = _csv.writer(f); w.writerow(["key","value"])
                if isinstance(data,dict):
                    for k,v in data.items(): w.writerow([k, json.dumps(v) if isinstance(v,(dict,list)) else str(v)])
        await cb(95,"Done!")
        return self._result(odir, fn)

    async def _json_txt(self, inp, tgt, odir, cb, stem):
        await cb(30, "JSON → TXT…")
        import json
        with open(inp,encoding="utf-8",errors="replace") as f: data = json.load(f)
        fn = f"{stem}.txt"
        self._write(odir, fn, json.dumps(data, ensure_ascii=False, indent=2))
        await cb(95,"Done!")
        return self._result(odir, fn)

    async def _json_html(self, inp, tgt, odir, cb, stem):
        await cb(30, "JSON → HTML…")
        import json
        with open(inp,encoding="utf-8",errors="replace") as f: data = json.load(f)
        if isinstance(data,list) and data and isinstance(data[0],dict):
            headers = list(data[0].keys())
            header_html = "".join(f"<th>{_html.escape(str(h))}</th>" for h in headers)
            rows_html = "".join(
                "<tr>"+"".join(f"<td>{_html.escape(str(row.get(h,'')))}</td>" for h in headers)+"</tr>"
                for row in data)
            body = f"<table><thead><tr>{header_html}</tr></thead><tbody>{rows_html}</tbody></table>"
        else:
            body = f"<pre>{_html.escape(json.dumps(data,ensure_ascii=False,indent=2))}</pre>"
        fn = f"{stem}.html"
        self._write(odir, fn, self._base_html(stem, body)); await cb(95,"Done!")
        return self._result(odir, fn)

    async def _json_xlsx(self, inp, tgt, odir, cb, stem):
        await cb(30, "JSON → XLSX…")
        import json
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill
        except ImportError:
            raise RuntimeError("Run: pip install openpyxl")
        with open(inp,encoding="utf-8",errors="replace") as f: data = json.load(f)
        wb = openpyxl.Workbook(); ws = wb.active; ws.title = stem[:31]
        if isinstance(data,list) and data and isinstance(data[0],dict):
            headers = list(data[0].keys()); ws.append(headers)
            for cell in ws[1]:
                cell.font = Font(bold=True)
                cell.fill = PatternFill("solid", fgColor="E8E8F0")
            for row in data: ws.append([str(row.get(h,"")) for h in headers])
        else:
            ws.append(["Data"]); ws.append([json.dumps(data,ensure_ascii=False)])
        fn = f"{stem}.xlsx"; wb.save(str(odir/fn)); await cb(95,"Done!")
        return self._result(odir, fn)

    async def _json_xml(self, inp, tgt, odir, cb, stem):
        await cb(30, "JSON → XML…")
        import json
        def to_xml(data, tag="item", indent=0):
            pad = "  "*indent
            if isinstance(data,dict):
                inner = "".join(to_xml(v, _safe_tag(k), indent+1) for k,v in data.items())
                return f"{pad}<{tag}>\n{inner}{pad}</{tag}>\n"
            elif isinstance(data,list):
                return "".join(to_xml(item,"item",indent) for item in data)
            else:
                return f"{pad}<{tag}>{_html.escape(str(data))}</{tag}>\n"
        with open(inp,encoding="utf-8",errors="replace") as f: data = json.load(f)
        xml = '<?xml version="1.0" encoding="UTF-8"?>\n' + to_xml(data,"root")
        fn = f"{stem}.xml"
        self._write(odir, fn, xml); await cb(95,"Done!")
        return self._result(odir, fn)

    # ── XML conversions ───────────────────────────────────────────────────────
    async def _xml_json(self, inp, tgt, odir, cb, stem):
        await cb(30, "XML → JSON…")
        import json, xml.etree.ElementTree as ET
        def elem_to_dict(el):
            d = {**el.attrib}
            for child in el:
                key = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                val = elem_to_dict(child) if len(child) else (child.text or "")
                if key in d:
                    if not isinstance(d[key],list): d[key]=[d[key]]
                    d[key].append(val)
                else: d[key]=val
            if el.text and el.text.strip() and not len(el): return el.text.strip()
            return d
        tree = ET.parse(str(inp)); root = tree.getroot()
        tag = root.tag.split("}")[-1] if "}" in root.tag else root.tag
        result = {tag: elem_to_dict(root)}
        fn = f"{stem}.json"
        (odir/fn).write_text(json.dumps(result,ensure_ascii=False,indent=2),encoding="utf-8")
        await cb(95,"Done!")
        return self._result(odir, fn)

    async def _xml_txt(self, inp, tgt, odir, cb, stem):
        await cb(30, "XML → TXT…")
        import xml.etree.ElementTree as ET
        tree = ET.parse(str(inp))
        parts = [el.text.strip() for el in tree.iter() if el.text and el.text.strip()]
        fn = f"{stem}.txt"
        self._write(odir, fn, "\n".join(parts)); await cb(95,"Done!")
        return self._result(odir, fn)

    async def _xml_html(self, inp, tgt, odir, cb, stem):
        await cb(30, "XML → HTML…")
        text = inp.read_text(encoding="utf-8", errors="replace")
        fn = f"{stem}.html"
        self._write(odir, fn, self._base_html(stem, f"<pre>{_html.escape(text)}</pre>"))
        await cb(95,"Done!")
        return self._result(odir, fn)

    async def _xml_pdf(self, inp, tgt, odir, cb, stem):
        await cb(30, "XML → PDF…")
        text = inp.read_text(encoding="utf-8", errors="replace")
        # Write as TXT then convert
        tmp = odir / f"_tmp_{stem}.txt"
        tmp.write_text(text, encoding="utf-8")
        try:
            r = await self._txt_pdf(tmp, "pdf", odir, cb, stem)
        finally:
            tmp.unlink(missing_ok=True)
        return r

    # ── CSS ───────────────────────────────────────────────────────────────────
    async def _copy_txt(self, inp, tgt, odir, cb, stem):
        await cb(30, "Reading file…")
        fn = f"{stem}.txt"
        self._write(odir, fn, inp.read_text(encoding="utf-8-sig", errors="replace"))
        await cb(95,"Done!")
        return self._result(odir, fn)

    async def _css_html(self, inp, tgt, odir, cb, stem):
        await cb(30, "CSS → HTML viewer…")
        css = _html.escape(inp.read_text(encoding="utf-8", errors="replace"))
        body = f'<h2 style="font-family:monospace;color:#555">📄 {_html.escape(inp.name)}</h2><pre><code class="css">{css}</code></pre>'
        fn = f"{stem}.html"
        self._write(odir, fn, self._base_html(stem, body)); await cb(95,"Done!")
        return self._result(odir, fn)

    # ── PPT (old .ppt format) handler ───────────────────────────────────────
    async def _ppt_convert(self, inp, tgt, odir, cb, stem):
        """Convert legacy .ppt using LibreOffice headless → .pptx → target."""
        await cb(15, "Converting .ppt via LibreOffice...")
        import shutil as _shutil
        soffice = _shutil.which("soffice") or _shutil.which("libreoffice")
        if not soffice:
            raise RuntimeError(
                ".ppt (old PowerPoint) requires LibreOffice:\n"
                "  Ubuntu: sudo apt install libreoffice\n"
                "  macOS:  brew install --cask libreoffice\n"
                "Or save as .pptx in PowerPoint first."
            )
        # Step 1: LibreOffice converts .ppt → .pptx in odir
        await cb(25, "LibreOffice: .ppt → .pptx...")
        # Run LibreOffice in executor (blocking) with subprocess.run
        import subprocess as _sp, uuid as _uuid, shutil as _sh
        tmp_profile = f"/tmp/lo_{_uuid.uuid4().hex[:8]}"
        def _run_lo():
            try:
                return _sp.run(
                    [soffice,
                     f"-env:UserInstallation=file://{tmp_profile}",
                     "--headless", "--convert-to", "pptx",
                     "--outdir", str(odir), str(inp)],
                    capture_output=True, timeout=90
                )
            finally:
                _sh.rmtree(tmp_profile, ignore_errors=True)
        loop = asyncio.get_event_loop()
        result = await loop.run_in_executor(None, _run_lo)
        pptx_path = odir / f"{inp.stem}.pptx"
        if result.returncode != 0 or not pptx_path.exists():
            raise RuntimeError(f"LibreOffice failed: {result.stderr.decode(errors='replace')[:200]}")
        if tgt == "pptx":
            await cb(95, "Done!")
            return self._result(odir, f"{stem}.pptx") if stem != inp.stem else self._result(odir, f"{inp.stem}.pptx")
        # Step 2: Now convert the .pptx to the actual target
        await cb(50, f"Converting .pptx → .{tgt}...")
        result = await self.convert(pptx_path, tgt, odir, cb, stem)
        return result


    # ── Pandoc universal fallback ─────────────────────────────────────────────
    async def _pan(self, inp, tgt, odir, cb, stem):
        # Guard: Pandoc cannot read these formats — route through our bridges instead
        ext = inp.suffix.lstrip(".").lower()
        if ext == "pdf":
            raise RuntimeError(
                f"Pandoc cannot read PDF files. Use the dedicated pdf→txt→{tgt} bridge."
            )
        if ext == "pptx":
            raise RuntimeError(
                f"Pandoc cannot read PPTX as input. Use the pptx→txt→{tgt} bridge."
            )
        await cb(30, "Converting with Pandoc…")
        pan = _find_pandoc()
        if not pan:
            raise RuntimeError(
                "Pandoc is not installed!\n"
                "  macOS:   brew install pandoc\n"
                "  Ubuntu:  sudo apt install pandoc\n"
                "  Windows: winget install pandoc\n"
                "  All:     pip install pandoc\n"
                "Download: https://pandoc.org/installing.html"
            )
        fn = f"{stem}.{tgt}"
        out = odir / fn
        cmd = [pan, str(inp), "-o", str(out), "--standalone", "--quiet"]
        if tgt == "pdf": cmd += ["--pdf-engine=xelatex"]
        # For EPUB, include title
        if tgt == "epub": cmd += [f"--metadata=title:{_html.escape(stem)}"]
        proc = await asyncio.create_subprocess_exec(
            *cmd, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
        try:
            _, err = await asyncio.wait_for(proc.communicate(), timeout=12.0)
        except asyncio.TimeoutError:
            proc.kill()
            raise RuntimeError(f"Pandoc timed out (12s limit). File may be too complex.")
        if proc.returncode != 0:
            err_s = err.decode(errors="replace")
            if tgt == "pdf" and ("xelatex" in err_s or "not found" in err_s.lower()):
                cmd2 = [c for c in cmd if c != "--pdf-engine=xelatex"]
                p2 = await asyncio.create_subprocess_exec(
                    *cmd2, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE)
                _, e2 = await p2.communicate()
                if p2.returncode != 0:
                    raise RuntimeError(
                        f"Pandoc PDF failed. Install LaTeX:\n"
                        f"  macOS: brew install --cask basictex\n"
                        f"  Ubuntu: sudo apt install texlive-xetex\n"
                        f"Error: {e2.decode(errors='replace')[-400:]}"
                    )
            else:
                raise RuntimeError(f"Pandoc failed:\n{err_s[-500:]}")
        await cb(95,"Done!")
        return self._result(odir, fn)


    # ═══════════════════════════════════════════════════════════════
    # RC-1: pdf → any  (2-step: pdf→txt→target)
    # ═══════════════════════════════════════════════════════════════
    async def _pdf_via_txt(self, inp, tgt, odir, cb, stem):
        """Route PDF to formats Pandoc can't read by going pdf→txt→target."""
        await cb(10, f"PDF → TXT → {tgt.upper()}…")
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            txt_res = await self._pdf_txt(inp, "txt", tmp_path, cb, stem)
            txt_file = tmp_path / txt_res["filename"]
            await cb(55, f"TXT → {tgt.upper()}…")
            return await self.convert(txt_file, tgt, odir, cb, stem)

    # ═══════════════════════════════════════════════════════════════
    # RC-2: pptx → any  (2-step: pptx→txt→target via python-pptx)
    # ═══════════════════════════════════════════════════════════════
    async def _pptx_via_txt(self, inp, tgt, odir, cb, stem):
        """Route PPTX to formats Pandoc can't read by extracting text first."""
        await cb(10, f"PPTX → TXT → {tgt.upper()}…")
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            txt_res = await self._pptx_txt(inp, "txt", tmp_path, cb, stem)
            txt_file = tmp_path / txt_res["filename"]
            await cb(55, f"TXT → {tgt.upper()}…")
            return await self.convert(txt_file, tgt, odir, cb, stem)

    # ═══════════════════════════════════════════════════════════════
    # RC-4: doc/odt/rtf/epub/json → pdf  (html bridge, no LaTeX)
    # ═══════════════════════════════════════════════════════════════
    async def _doc_via_html_pdf(self, inp, tgt, odir, cb, stem):
        """Convert any document to PDF via HTML bridge (WeasyPrint, no LaTeX)."""
        await cb(10, f"{inp.suffix.upper()} → HTML → PDF…")
        ext = inp.suffix.lstrip(".").lower()
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            tmp_path = Path(tmp)
            # Step 1: to HTML
            if ext == "json":
                html_res = await self._json_html(inp, "html", tmp_path, cb, stem)
            elif ext in ("odt", "rtf", "epub"):
                html_res = await self._pan(inp, "html", tmp_path, cb, stem)
            else:
                html_res = await self._pan(inp, "html", tmp_path, cb, stem)
            html_file = tmp_path / html_res["filename"]
            await cb(60, "HTML → PDF…")
            # Step 2: HTML → PDF via WeasyPrint
            return await self._html_pdf(html_file, "pdf", odir, cb, stem)

    # ═══════════════════════════════════════════════════════════════
    # RC-5: text/document → CSV  (native writer, no Pandoc template)
    # ═══════════════════════════════════════════════════════════════
    async def _txt_csv(self, inp, tgt, odir, cb, stem):
        """Convert text-based formats to CSV: one column per non-empty line."""
        await cb(30, f"{inp.suffix.upper()} → CSV…")
        ext = inp.suffix.lstrip(".").lower()
        # For XML, extract text nodes
        if ext == "xml":
            import xml.etree.ElementTree as ET
            try:
                tree = ET.parse(str(inp))
                lines = [el.text.strip() for el in tree.iter()
                         if el.text and el.text.strip()]
            except Exception:
                lines = inp.read_text(encoding="utf-8", errors="replace").splitlines()
        else:
            lines = inp.read_text(encoding="utf-8-sig", errors="replace").splitlines()
        fn = f"{stem}.csv"
        with open(odir / fn, "w", newline="", encoding="utf-8") as f:
            w = _csv.writer(f)
            w.writerow(["line", "content"])
            for i, line in enumerate(lines, 1):
                stripped = line.strip()
                if stripped:
                    w.writerow([i, stripped])
        await cb(95, "Done!")
        return self._result(odir, fn)

    async def _html_csv(self, inp, tgt, odir, cb, stem):
        """HTML → CSV: extract all table data; fallback to text lines."""
        await cb(30, "HTML → CSV…")
        from html.parser import HTMLParser
        class TableParser(HTMLParser):
            def __init__(self):
                super().__init__()
                self.tables, self._cur, self._row, self._cell = [], [], [], []
                self._in_cell = False
            def handle_starttag(self, tag, attrs):
                if tag == "tr": self._row = []
                elif tag in ("td","th"): self._cell = []; self._in_cell = True
            def handle_endtag(self, tag):
                if tag in ("td","th"):
                    self._row.append("".join(self._cell).strip())
                    self._in_cell = False
                elif tag == "tr" and self._row:
                    self._cur.append(self._row); self._row = []
                elif tag == "table" and self._cur:
                    self.tables.append(self._cur); self._cur = []
            def handle_data(self, data):
                if self._in_cell: self._cell.append(data)
        p = TableParser()
        p.feed(inp.read_text(encoding="utf-8", errors="replace"))
        fn = f"{stem}.csv"
        with open(odir / fn, "w", newline="", encoding="utf-8") as f:
            w = _csv.writer(f)
            if p.tables:
                for tbl in p.tables:
                    for row in tbl: w.writerow(row)
                    w.writerow([])
            else:
                # No tables — write text lines
                lines = inp.read_text(encoding="utf-8", errors="replace").splitlines()
                w.writerow(["line","content"])
                for i, l in enumerate(lines, 1):
                    if l.strip(): w.writerow([i, l.strip()])
        await cb(95, "Done!")
        return self._result(odir, fn)

    async def _docx_csv(self, inp, tgt, odir, cb, stem):
        """DOCX → CSV: extract tables first, then paragraphs as fallback."""
        await cb(30, "DOCX → CSV…")
        from docx import Document
        d = Document(str(inp))
        fn = f"{stem}.csv"
        with open(odir / fn, "w", newline="", encoding="utf-8") as f:
            w = _csv.writer(f)
            if d.tables:
                for tbl in d.tables:
                    for row in tbl.rows:
                        w.writerow([c.text.strip() for c in row.cells])
                    w.writerow([])
            else:
                w.writerow(["line", "content"])
                for i, p in enumerate([p for p in d.paragraphs if p.text.strip()], 1):
                    w.writerow([i, p.text.strip()])
        await cb(95, "Done!")
        return self._result(odir, fn)

    # ═══════════════════════════════════════════════════════════════
    # RC-3: text/doc → xlsx  (real openpyxl XLSX, not Pandoc fake)
    # ═══════════════════════════════════════════════════════════════
    async def _txt_xlsx(self, inp, tgt, odir, cb, stem):
        """Convert text/markup to real XLSX via openpyxl (no Pandoc fake xlsx)."""
        await cb(30, f"{inp.suffix.upper()} → XLSX…")
        import openpyxl
        from openpyxl.styles import Font, PatternFill
        ext = inp.suffix.lstrip(".").lower()
        wb = openpyxl.Workbook(); ws = wb.active; ws.title = stem[:31]
        if ext == "xml":
            import xml.etree.ElementTree as ET
            try:
                tree = ET.parse(str(inp))
                lines = [(el.tag.split("}")[-1], el.text or "")
                         for el in tree.iter() if el.text and el.text.strip()]
                ws.append(["Tag", "Value"])
                for tag, val in lines: ws.append([tag, val.strip()])
            except Exception:
                lines = inp.read_text(encoding="utf-8", errors="replace").splitlines()
                ws.append(["Line", "Content"])
                for i, l in enumerate(lines, 1):
                    if l.strip(): ws.append([i, l.strip()])
        else:
            lines = inp.read_text(encoding="utf-8-sig", errors="replace").splitlines()
            ws.append(["Line", "Content"])
            for i, l in enumerate(lines, 1):
                if l.strip(): ws.append([i, l.strip()])
        # Style header row
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="E8E8F0")
        fn = f"{stem}.xlsx"; wb.save(str(odir / fn))
        await cb(95, "Done!")
        return self._result(odir, fn)


    async def _xlsx_via_txt(self, inp, tgt, odir, cb, stem):
        """XLSX → {docx,epub,md,odt,rtf,pptx} via txt intermediate."""
        import tempfile
        await cb(10, f"XLSX → TXT → {tgt.upper()}…")
        with tempfile.TemporaryDirectory() as tmp:
            txt_res = await self._xlsx_txt(inp, "txt", Path(tmp), cb, stem)
            txt_file = Path(tmp) / txt_res["filename"]
            await cb(55, f"TXT → {tgt.upper()}…")
            return await self.convert(txt_file, tgt, odir, cb, stem)

    async def _doc_via_txt_xlsx(self, inp, tgt, odir, cb, stem):
        """ODT/EPUB → XLSX via pandoc→txt→openpyxl (avoids reading ZIP as text)."""
        import tempfile
        await cb(10, f"{inp.suffix.upper()} → TXT → XLSX…")
        with tempfile.TemporaryDirectory() as tmp:
            # Use Pandoc to get clean text from ODT/EPUB
            txt_res = await self._pan(inp, "txt", Path(tmp), cb, stem)
            txt_file = Path(tmp) / txt_res["filename"]
            await cb(55, "TXT → XLSX…")
            return await self._txt_xlsx(txt_file, "xlsx", odir, cb, stem)


    async def _docx_xlsx(self, inp, tgt, odir, cb, stem):
        """DOCX → real XLSX via openpyxl: tables first, then paragraphs."""
        await cb(30, "DOCX → XLSX…")
        import openpyxl
        from openpyxl.styles import Font, PatternFill
        from docx import Document
        wb = openpyxl.Workbook(); ws = wb.active; ws.title = stem[:31]
        d = Document(str(inp)); row_num = 1
        if d.tables:
            for ti, tbl in enumerate(d.tables):
                if ti > 0:
                    ws.cell(row=row_num, column=1, value="─" * 20); row_num += 1
                for r in tbl.rows:
                    for ci, cell in enumerate(r.cells, 1):
                        c = ws.cell(row=row_num, column=ci, value=cell.text.strip())
                        if row_num == 1: c.font = Font(bold=True)
                    row_num += 1
        if not d.tables:
            ws.append(["#", "Content"]); row_num += 1
            for p in d.paragraphs:
                if p.text.strip():
                    ws.append([row_num - 1, p.text.strip()]); row_num += 1
        for cell in ws[1]:
            cell.font = Font(bold=True)
            cell.fill = PatternFill("solid", fgColor="DCE6F1")
        fn = f"{stem}.xlsx"; wb.save(str(odir / fn))
        await cb(95, "Done!"); return self._result(odir, fn)

    async def _csv_docx(self, inp, tgt, odir, cb, stem):
        """CSV → DOCX: proper Word table with styled header row."""
        await cb(30, "CSV → DOCX…")
        import csv as _csvm
        from docx import Document
        from docx.shared import Pt
        d = Document(); d.core_properties.title = stem
        rows = []
        with open(inp, "r", encoding="utf-8-sig", errors="replace") as f:
            rows = list(_csvm.reader(f))
        if not rows:
            d.add_paragraph("(empty CSV)")
        else:
            cols = max(len(r) for r in rows)
            tbl = d.add_table(rows=len(rows), cols=cols)
            tbl.style = "Table Grid"
            for ri, row in enumerate(rows):
                for ci in range(cols):
                    val = row[ci] if ci < len(row) else ""
                    cell = tbl.cell(ri, ci)
                    p = cell.paragraphs[0]
                    run = p.add_run(val)
                    if ri == 0: run.bold = True
        fn = f"{stem}.docx"; d.save(str(odir / fn))
        await cb(95, "Done!"); return self._result(odir, fn)

    async def _csv_md(self, inp, tgt, odir, cb, stem):
        """CSV → Markdown table."""
        await cb(30, "CSV → MD…")
        import csv as _csvm
        rows = []
        with open(inp, "r", encoding="utf-8-sig", errors="replace") as f:
            rows = list(_csvm.reader(f))
        lines = []
        if rows:
            lines.append("| " + " | ".join(rows[0]) + " |")
            lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
            for row in rows[1:]:
                lines.append("| " + " | ".join(row) + " |")
        fn = f"{stem}.md"
        self._write(odir, fn, "\n".join(lines))
        await cb(95, "Done!"); return self._result(odir, fn)


def _safe_tag(name: str) -> str:
    """Make a string safe to use as XML tag name."""
    tag = re.sub(r'[^a-zA-Z0-9_\-]', '_', str(name))
    if tag and tag[0].isdigit(): tag = '_' + tag
    return tag or 'field'
